import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
import matplotlib.pyplot as plt
import numpy as np

def create_tectonic_map():
    fig, ax = plt.subplots(figsize=(5,2.5))
    # World map as background
    ax.set_xlim(-180, 180)
    ax.set_ylim(-60, 90)
    ax.set_xticks([])
    ax.set_yticks([])
    ax.axis('off')

    # Draw tectonic plates as circles
    ax.scatter([-100], [40], s=1500, c='dodgerblue', alpha=0.35, label='US')         # US Plate
    ax.scatter([105], [35], s=1800, c='red', alpha=0.28, label='China')              # China Plate
    ax.scatter([10], [50], s=1500, c='limegreen', alpha=0.32, label='Western Europe')# EU Plate

    # Add annotation labels
    ax.annotate('US', xy=(-100, 55), fontsize=13, weight='bold', color='blue')
    ax.annotate('China', xy=(110, 50), fontsize=13, weight='bold', color='red')
    ax.annotate('Western Europe', xy=(10, 60), fontsize=13, weight='bold', color='green')
    ax.set_title("Wilo's Three Tectonic Plates: Strategic Centers", fontsize=16, pad=18)
    buf = BytesIO()
    plt.tight_layout()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True)
    plt.close()
    buf.seek(0)
    return buf

def create_swot_chart():
    fig, ax = plt.subplots(figsize=(6,3.2))
    ax.axis('off')
    # Draw a 2x2 grid for SWOT
    swot = [
        ['Strengths', 'Weaknesses'],
        ['Opportunities', 'Threats']
    ]
    cell_text = [
        [
            "• Multi-regional presence\n• Supply chain resilience\n• Localized R&D",
            "• Increased complexity\n• Potential for silos/duplication"
        ],
        [
            "• Growth in Asia/US\n• Innovation hubs\n• Talent acquisition",
            "• Geopolitical risks\n• Regulatory uncertainty\n• Integration challenges"
        ]
    ]
    table = ax.table(
        cellText=cell_text,
        colLabels=["Internal", "Internal"],
        rowLabels=["Strengths", "Opportunities"],
        loc='center',
        cellLoc='left'
    )
    table.auto_set_font_size(False)
    table.set_fontsize(11)
    for (row, col), cell in table.get_celld().items():
        cell.set_height(0.18)
        if row == 0:
            cell.set_facecolor('#dbefff')
        elif row == 1:
            cell.set_facecolor('#eaffdb')
    ax.set_title("SWOT: Tri-Polar HQ Strategy", fontsize=15, pad=16)
    buf = BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True)
    plt.close()
    buf.seek(0)
    return buf

def create_ir_framework():
    fig, ax = plt.subplots(figsize=(5,3))
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_xticks([])
    ax.set_yticks([])
    # Draw axes
    ax.arrow(0.05, 0.05, 0.85, 0, head_width=0.02, head_length=0.03, fc='black', ec='black')
    ax.arrow(0.05, 0.05, 0, 0.85, head_width=0.03, head_length=0.02, fc='black', ec='black')
    ax.text(0.48, 0.01, "Global Integration", ha='center', fontsize=12)
    ax.text(0.01, 0.48, "Local Responsiveness", va='center', rotation='vertical', fontsize=12)
    # Plot Wilo's strategy position
    ax.scatter([0.8], [0.7], s=250, c='orange')
    ax.text(0.82, 0.75, "Wilo", fontsize=12, fontweight='bold', color='orange')
    # Annotate quadrants
    ax.text(0.7, 0.9, "Transnational", fontsize=11, weight='bold', color='blue')
    ax.text(0.15, 0.88, "Multi-domestic", fontsize=10, color='grey')
    ax.text(0.72, 0.1, "Global", fontsize=10, color='grey')
    ax.text(0.14, 0.14, "International", fontsize=10, color='grey')
    ax.set_title("Integration–Responsiveness (I-R) Framework", fontsize=15)
    buf = BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True)
    plt.close()
    buf.seek(0)
    return buf

def create_leadership_radar():
    # Radar for 6 leadership qualities
    labels = [
        "Vision & Strategy", "Decisiveness", "Empowerment",
        "Cultural Intelligence", "Change Advocacy", "Systems Thinking"
    ]
    stats = [9, 8, 8, 9, 8, 8]  # arbitrary high scores for Wilo
    angles = np.linspace(0, 2 * np.pi, len(labels), endpoint=False).tolist()
    stats += stats[:1]  # repeat first value to close radar
    angles += angles[:1]
    fig, ax = plt.subplots(figsize=(4.5,4), subplot_kw=dict(polar=True))
    ax.plot(angles, stats, color="orangered", linewidth=2)
    ax.fill(angles, stats, color="orange", alpha=0.25)
    ax.set_yticks([])
    ax.set_xticks(angles[:-1])
    ax.set_xticklabels(labels, fontsize=11)
    ax.set_title("Leadership Qualities for Change", y=1.10, fontsize=14)
    buf = BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', transparent=True)
    plt.close()
    buf.seek(0)
    return buf

def add_image_slide(prs, title, img_stream, width=6.5, height=3.6):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    left = Inches(0.6)
    top = Inches(1.2)
    pic = slide.shapes.add_picture(img_stream, left, top, width=Inches(width), height=Inches(height))

def create_wilo_ppt():
    prs = Presentation()
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Wilo SE: Managing Organizational Change in Times of Deglobalization"
    slide.placeholders[1].text = "Change Management Case Analysis\nMBA Report Summary"

    # Slide 1: Introduction
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Introduction"
    slide.placeholders[1].text = (
        "• Wilo SE: global German-based manufacturer of pumps & systems.\n"
        "• CEO Hermes: diagnosed 'three tectonic plates': US, China, Western Europe.\n"
        "• Reflects deglobalization, geopolitical fragmentation, regionalization.\n"
        "• Adopted a tri-polar HQ model."
    )

    # Slide 2: Tectonic Plates Visual
    tectonic_img = create_tectonic_map()
    add_image_slide(prs, "The 'Three Tectonic Plates' (US, China, Europe)", tectonic_img)

    # Slide 3: Do you agree with the metaphor?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Agreement with 'Three Tectonic Plates' Metaphor"
    tf = slide.placeholders[1].text_frame
    tf.text = "Yes – Realistically captures today’s fragmented global environment."
    for bullet in [
        "• Geopolitical fragmentation: protectionism, divergent regulation.",
        "• Proximity to customers, regulators, talent is crucial.",
        "• Multi-HQ = resilience, agility, market relevance.",
        "• Supported by VUCA, Lewin’s, Kotter’s, Bartlett & Ghoshal."
    ]:
        tf.add_paragraph().text = bullet

    # Slide 4: I-R Framework Visual
    ir_img = create_ir_framework()
    add_image_slide(prs, "Wilo's Strategic Position: I-R Framework", ir_img, width=5.5, height=3)

    # Slide 5: Why China HQ?
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Why Establish a Second HQ in China?"
    tf = slide.placeholders[1].text_frame
    tf.text = "Beyond geopolitics, key drivers include:"
    for bullet in [
        "• Access to fastest-growing market (infrastructure, digital projects).",
        "• Tap into China’s digital innovation ecosystem.",
        "• Attract & retain local talent.",
        "• Customer proximity, customization, fast delivery.",
        "• Strengthen supply chain resilience."
    ]:
        tf.add_paragraph().text = bullet

    # Slide 6: SWOT Visual
    swot_img = create_swot_chart()
    add_image_slide(prs, "SWOT: Tri-Polar HQ Strategy", swot_img)

    # Slide 7: Organizational Implications
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Organizational Implications of Multi-HQ"
    tf = slide.placeholders[1].text_frame
    tf.text = "Profound changes for structure, systems, and culture:"
    for bullet in [
        "• Decentralized decision-making, regional empowerment.",
        "• Hybrid reporting: regional innovation, global oversight.",
        "• Global governance: regional heads on board.",
        "• Cross-cultural talent management.",
        "• Risk of silos/duplication (needs active management)."
    ]:
        tf.add_paragraph().text = bullet

    # Slide 8: Leadership Qualities Visual
    radar_img = create_leadership_radar()
    add_image_slide(prs, "Leadership Qualities for Change", radar_img, width=4.6, height=3.5)

    # Slide 9: Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion"
    tf = slide.placeholders[1].text_frame
    tf.text = "Tri-polar HQ strategy = robust, proactive response to deglobalization."
    for bullet in [
        "• 'Three tectonic plates' approach: analytic and actionable.",
        "• Balances global cohesion with local empowerment.",
        "• Success: visionary, collaborative, and culturally agile leadership.",
        "• A blueprint for resilient multinationals."
    ]:
        tf.add_paragraph().text = bullet

    # Slide 10: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Thank You!"
    slide.placeholders[1].text = "Questions?\nWilo SE: Managing Organizational Change in Times of Deglobalization"

    pptx_io = BytesIO()
    prs.save(pptx_io)
    pptx_io.seek(0)
    return pptx_io

# Streamlit UI
st.set_page_config(page_title="Wilo SE Change Management PPT Generator", layout="centered")
st.title("Wilo SE: Organizational Change in Times of Deglobalization")
st.write("""
Generate a consulting-style PowerPoint summarizing the Wilo SE case with key frameworks and visuals.
""")

if st.button("Generate & Download PPTX"):
    pptx_io = create_wilo_ppt()
    st.download_button(
        label="Download Wilo_SE_Change_Management.pptx",
        data=pptx_io,
        file_name="Wilo_SE_Change_Management.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
    )
    st.success("PPTX ready! Download using the button above.")

st.markdown("""
---
#### Slide Visuals Include:
- Tectonic plates world map
- SWOT chart (Tri-Polar HQ strategy)
- Integration-Responsiveness (I-R) framework
- Radar chart for leadership qualities

**Just download and present!**  
Customize in PowerPoint as needed.

---
""")
